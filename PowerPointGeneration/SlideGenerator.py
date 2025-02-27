import streamlit as st
import requests
import os
import random
import base64
from io import BytesIO
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Pt, Inches
from PIL import Image

#  LangChain & Model Management
from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate
from Model_Manager import ModelManager

#  Configurations & Base Structure
from BaseApp import BaseApp  # Importing shared UI structure
from Config import PRIMARY_MODEL, SECONDARY_MODEL

load_dotenv()
API_KEY = os.getenv("SEGMIND_API_KEY")

#  Function to find the correct text placeholder
def find_text_placeholder(slide):
    """Finds the first available text placeholder in a slide."""
    for shape in slide.shapes:
        if shape.is_placeholder and ("Content Placeholder" in shape.name or "Text Placeholder" in shape.name):
            return shape.text_frame
    return None  # No text placeholder found

#  Function to find the correct image placeholder
def find_image_placeholder(slide):
    """Finds the first available image placeholder in a slide."""
    for shape in slide.shapes:
        if shape.is_placeholder and "Picture Placeholder" in shape.name:
            return shape
    return None  # No image placeholder found

def generate_image_stable_diffusion(prompt, negative_prompt="low quality, blurry", width=904, height=1260):
    """
    Generates an image using Stable Diffusion with the specified prompt and dimensions.

    Args:
        prompt (str): The description of the image to generate.
        negative_prompt (str): Elements to avoid in the image.
        width (int): Width of the generated image in pixels.
        height (int): Height of the generated image in pixels.

    Returns:
        Image or None: The generated image or None if failed.
    """
    url = "https://api.segmind.com/v1/stable-diffusion-3.5-turbo-txt2img"
    headers = {"x-api-key": API_KEY}

    random_seed = random.randint(0, 99999999999)
    data = {
        "prompt": prompt,
        "negative_prompt": negative_prompt,
        "steps": 4,
        "guidance_scale": 1,
        "seed": random_seed,
        "sampler": "dpmpp_2m",
        "scheduler": "sgm_uniform",
        "width": width,
        "height": height,
        "aspect_ratio": "custom",
        "batch_size": 1,
        "image_format": "jpeg",
        "image_quality": 95,
        "base64": True
    }

    try:
        response = requests.post(url, json=data, headers=headers)
        response.raise_for_status()
        response_data = response.json()

        image_base64 = response_data.get("image", "")
        if image_base64:
            img_data = base64.b64decode(image_base64)
            img = Image.open(BytesIO(img_data))
            return img
        else:
            return None
    except Exception as e:
        print(f"Stable Diffusion Error: {e}")
        return None
# 🔥 Define Templates Path
TEMPLATE_DIR = "./Pres_templates"
available_templates = ["Celestial", "Circuit", "Ion", "Mesh", "Retrospect", "Atlas", "Wisp", "Gallery", "Madison", "Organic"]

class SlideGenerator(BaseApp):
    def __init__(self, model_manager, app_name="Slide Generator", app_slogan="📊 Create AI-Powered Presentations!"):
        super().__init__(app_name)
        self.app_slogan = app_slogan
        self.model_manager = model_manager
        if "history" not in st.session_state:
            st.session_state.history = []

        # Initialize session state variables
        if "user_info" not in st.session_state:
            st.session_state.user_info = {
                "num_slides": 10,
                "presentation_request": "",
                "document_content": "",
            }
        if "generate_clicked" not in st.session_state:
            st.session_state.generate_clicked = False

    def welcome_screen(self):
        super().welcome_screen(self.app_slogan)

    def display_side(self):
        """Sidebar for user to specify slide preferences."""
        with st.sidebar:
            st.header("⚙️ Slide Settings")

            # 🔥 Store selected template in session state & theme name
            st.session_state.theme = st.selectbox(
                "🎨 Select a PowerPoint Template", available_templates, index=0
            )

            # 🔥 Set the template path based on selection
            st.session_state.final_template = os.path.join(TEMPLATE_DIR, f"{st.session_state.theme}.pptx")

            # 🔥 Dropdown for number of slides
            st.session_state.num_slides = st.selectbox("📊 Number of Slides", [5, 10, 15, 20], index=1)

            # 🔥 Text input for user description
            st.session_state.presentation_request = st.text_area(
                "📌 Describe your presentation",
                placeholder="E.g., 'A professional presentation on AI in business with trends & ethical concerns.'"
            )

            if st.button("🎨 Generate Slides"):
                st.session_state.generate_clicked = True

    def generate_prompt(self):
        return ChatPromptTemplate.from_template("""
        You are an AI that generates structured PowerPoint slide content.
        When needed, suggest a detailed image description that can be used to generate an image and add it to the slide. The number of images should not exceed half the number of slides.
        If the slide does not require an image, return '#Image: None'.
        Follow this exact format but for the number of slides given

        #Title: TITLE OF THE PRESENTATION
        #Subtitle: A brief one-line description of the presentation

        #Slide: 1
        #Header: Table of Contents
        #Contents: This presentation includes:
        - Main topic 1
        - Main topic 2
        - Main topic 3
        #Image: None

        #Slide: 2
        #Header: TITLE OF SLIDE
        #Content: CONTENT OF THE SLIDE
        #Image: SUGGESTED IMAGE description OR None

        #Slide: 3
        #Header: TITLE OF SLIDE
        #Contents: Key points:
        - Point 1
          - Sub-point A
          - Sub-point B
        - Point 2
        #Image: SUGGESTED IMAGE description OR None

        Ensure that all slides follow this format exactly. Do not add extra slides.

        number of slides {num_slides}
        subject {presentation_request}
        """).format_prompt(
            num_slides=st.session_state.num_slides,
            presentation_request=st.session_state.presentation_request
        )



    def generate_ppt(self, ai_response):
        """Creates a PowerPoint presentation using AI response."""
        try:
            content = ai_response.content
            st.session_state.history.append(content)

            # 🔥 Debug: Store AI Response in Session History
            lines = content.split("\n")

            # 🔥 Extract Theme, Title, and Subtitle from LLM Response
            theme = st.session_state.theme  # Directly use selected theme
            title = None
            subtitle = None

            slide_numbers = []
            slide_data = {}  # Store slide data in a structured format

            for i, line in enumerate(lines):
                line = line.strip()
                if line.startswith("#Slide:"):
                    slide_num = int(line.replace("#Slide:", "").strip())
                    slide_numbers.append(slide_num)
                    slide_data[slide_num] = {
                        "title": "", "content": [], "image": None, "is_bullet_points": False
                    }
                elif line.startswith("#Title:"):
                    title = line.replace("#Title:", "").strip()
                elif line.startswith("#Subtitle:"):
                    subtitle = line.replace("#Subtitle:", "").strip()

            # 🔥 Load PowerPoint Template
            ppt_template = os.path.join(TEMPLATE_DIR, f"{theme}.pptx")
            prs = Presentation(ppt_template) if os.path.exists(ppt_template) else Presentation()

            # 🔥 Modify First Slide Instead of Adding a New One
            first_slide = prs.slides[0]  # Get the existing first slide
            first_slide.shapes.title.text = title  # Set the title

            # Modify subtitle if there's a placeholder
            for shape in first_slide.shapes:
                if shape.has_text_frame and shape.text_frame.text == "":
                    shape.text_frame.text = subtitle
                    break  # Stop after setting the first empty placeholder

            # 🔥 Process Lines & Store Slide Data
            current_slide = None
            for line in lines:
                line = line.strip()

                if line.startswith("#Slide:"):
                    current_slide = int(line.replace("#Slide:", "").strip())
                    continue

                elif current_slide is not None:
                    if line.startswith("#Header:"):
                        slide_data[current_slide]["title"] = line.replace("#Header:", "").strip()

                    elif line.startswith("#Contents:"):
                        slide_data[current_slide]["is_bullet_points"] = True
                        slide_data[current_slide]["content"].append(line.replace("#Contents:", "").strip())

                    elif line.startswith("#Content:"):
                        slide_data[current_slide]["is_bullet_points"] = False
                        slide_data[current_slide]["content"].append(line.replace("#Content:", "").strip())

                    elif line.startswith("- "):
                        slide_data[current_slide]["content"].append(line.replace("- ", "").strip())

                    elif line.startswith("#Image:"):
                        image_suggestion = line.replace("#Image:", "").strip()
                        slide_data[current_slide]["image"] = None if image_suggestion.lower() == "none" else image_suggestion

            # 🔥 Process Remaining Slides (Avoid Adding Empty Slides)
            slide_count = 0
            valid_slides = []  # Store non-empty slides

            for slide_num in sorted(slide_numbers):
                slide_info = slide_data[slide_num]

                # ✅ Skip Adding Slide if There's No Title or Content
                if not slide_info["title"] and not slide_info["content"]:
                    continue

                # 🔥 Choose Slide Layout (Fixed Since No Custom Templates)
                slide_layout_index = 8 if slide_info["image"] else 1
                slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
                slide.shapes.title.text = slide_info["title"]
                valid_slides.append(slide)

                # 🔥 Find Correct Text Placeholder
                text_placeholder = find_text_placeholder(slide)

                # 🔥 Insert Content
                if text_placeholder:
                    if slide_info["is_bullet_points"]:
                        for i, point in enumerate(slide_info["content"]):
                            p = text_placeholder.add_paragraph()
                            p.text = point.strip()
                            p.font.size = Pt(24) if i == 0 else Pt(18)
                            if i > 0:
                                p.level = 1
                    else:
                        text_placeholder.text = "\n".join(slide_info["content"])
                        text_placeholder.paragraphs[0].font.size = Pt(24)

                else:
                    # 🔥 If no content placeholder exists, move content into subtitle if available
                    subtitle_placeholder = None
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape.text_frame.text == "":
                            subtitle_placeholder = shape.text_frame
                            break

                    if subtitle_placeholder:
                        subtitle_placeholder.text = "\n".join(slide_info["content"])

                # 🔥 Generate & Insert Image
                if slide_info["image"]:
                    prompt = f"You are an AI designed to make pictures for PowerPoint presentations. Make a picture about {slide_info['image']}."
                    generated_image = generate_image_stable_diffusion(prompt)

                    if generated_image:
                        img_path = os.path.join("./images", f"{slide_num}.jpg")
                        generated_image.save(img_path, format="JPEG")

                        image_placeholder = find_image_placeholder(slide)

                        if image_placeholder:
                            image_placeholder.insert_picture(img_path)
                        else:
                            slide.shapes.add_picture(img_path, Inches(1), Inches(2), Inches(6), Inches(4))

                        os.remove(img_path)

            # 🔥 Cleanup: Remove Empty Slides Before Saving
            slides_to_remove = [s for s in prs.slides if not any(shape.has_text_frame and shape.text_frame.text.strip() for shape in s.shapes)]
            for s in slides_to_remove:
                prs.slides._sldIdLst.remove(s._element)  # Remove from XML tree

            # 🔥 Save PowerPoint
            ppt_filename = "Generated_Presentation.pptx"
            prs.save(ppt_filename)

            return ppt_filename if os.path.exists(ppt_filename) else None

        except Exception as e:
            st.error(f"Error generating PowerPoint: {e}")
            return None

    def handle_input(self):
        """Handles user input and generates a PowerPoint presentation."""

        # Ensure the generate button was clicked before running the process
        if "generate_clicked" in st.session_state and st.session_state.generate_clicked:
            if not st.session_state.presentation_request.strip():
                st.warning("Please describe the presentation you want.")
                return

            model = ChatOpenAI(temperature=0.2, model="gpt-4o-mini")

            # Generate response using LangChain
            conversation = self.build_conversation()
            response = model.invoke(conversation)

            # Store the response in session history
            if response:
                st.session_state.history.append(response.content)

            # Generate PowerPoint and save filename in session state
            ppt_filename = self.generate_ppt(response)
            if ppt_filename:
                st.session_state.ppt_filename = ppt_filename  # Store filename for later use
                st.session_state.generate_clicked = False  # Prevent re-triggering
                st.session_state.show_edit_button = True  # Ensure "Make Edits" remains visible

        # If edit mode is enabled, show the edit interface
        if "edit_mode" in st.session_state and st.session_state.edit_mode:
            user_edits = st.text_area("💬 Provide feedback for edits")

            if st.button("🔄 Submit Edits"):
                if st.session_state.history:
                    last_response = st.session_state.history[-1]
                    new_request = f"{last_response}\n\n# User Edits:\n{user_edits}\n\nEnsure the revised presentation follows the exact previous format."

                    # Send updated request to LLM
                    model = ChatOpenAI(temperature=0.2, model="gpt-4o-mini")
                    response = model.invoke([{"role": "user", "content": new_request}])

                    # Store the new response in session history
                    if response:
                        st.session_state.history.append(response.content)

                    # Generate new PPT with edits
                    ppt_filename = self.generate_ppt(response)

                    if ppt_filename:
                        st.session_state.ppt_filename = ppt_filename
                        st.session_state.edit_mode = False  # Exit edit mode
                        st.session_state.show_edit_button = True  # Ensure "Make Edits" remains visible

        # Ensure "Make Edits" button is always visible if a presentation exists
        if "ppt_filename" in st.session_state and st.session_state.ppt_filename:

            with open(st.session_state.ppt_filename, "rb") as f:
                st.download_button("📥 Download Presentation", f, st.session_state.ppt_filename,
                                   "application/vnd.openxmlformats-officedocument.presentationml.presentation")

            # Only show "Make Edits" if we're NOT in edit mode
            if "show_edit_button" in st.session_state and st.session_state.show_edit_button:
                if not st.session_state.get("edit_mode", False):  # Hides when editing
                    if st.button("📝 Make Edits"):
                        st.session_state.edit_mode = True  # Set immediately
                        st.rerun()  # 🔥 Force UI refresh to reflect changes immediately



# Run the app
if __name__ == "__main__":
    model_manager = ModelManager(PRIMARY_MODEL, SECONDARY_MODEL)
    app = SlideGenerator(model_manager)
    app.welcome_screen()
    app.display_side()
    app.handle_input()
