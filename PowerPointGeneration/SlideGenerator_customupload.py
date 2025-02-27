import streamlit as st
import requests
import os
import random
import base64
from io import BytesIO
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Pt, Inches
from PIL import Image

#  LangChain & Model Management
from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate
from Model_Manager import ModelManager

#  Configurations & Base Structure
from BaseApp import BaseApp  # Importing shared UI structure
from Config import PRIMARY_MODEL, SECONDARY_MODEL

load_dotenv()
API_KEY = os.getenv("SEGMIND_API_KEY")

#  Function to find the correct text placeholder
def find_text_placeholder(slide):
    """Finds the first available text placeholder in a slide."""
    for shape in slide.shapes:
        if shape.is_placeholder and ("Content Placeholder" in shape.name or "Text Placeholder" in shape.name):
            return shape.text_frame
    return None  # No text placeholder found

#  Function to find the correct image placeholder
def find_image_placeholder(slide):
    """Finds the first available image placeholder in a slide."""
    for shape in slide.shapes:
        if shape.is_placeholder and "Picture Placeholder" in shape.name:
            return shape
    return None  # No image placeholder found

def generate_image_stable_diffusion(prompt, negative_prompt="low quality, blurry", width=904, height=1260):
    """
    Generates an image using Stable Diffusion with the specified prompt and dimensions.

    Args:
        prompt (str): The description of the image to generate.
        negative_prompt (str): Elements to avoid in the image.
        width (int): Width of the generated image in pixels.
        height (int): Height of the generated image in pixels.

    Returns:
        Image or None: The generated image or None if failed.
    """
    url = "https://api.segmind.com/v1/stable-diffusion-3.5-turbo-txt2img"
    headers = {"x-api-key": API_KEY}

    random_seed = random.randint(0, 99999999999)
    data = {
        "prompt": prompt,
        "negative_prompt": negative_prompt,
        "steps": 4,
        "guidance_scale": 1,
        "seed": random_seed,
        "sampler": "dpmpp_2m",
        "scheduler": "sgm_uniform",
        "width": width,
        "height": height,
        "aspect_ratio": "custom",
        "batch_size": 1,
        "image_format": "jpeg",
        "image_quality": 95,
        "base64": True
    }

    try:
        response = requests.post(url, json=data, headers=headers)
        response.raise_for_status()
        response_data = response.json()

        image_base64 = response_data.get("image", "")
        if image_base64:
            img_data = base64.b64decode(image_base64)
            img = Image.open(BytesIO(img_data))
            return img
        else:
            return None
    except Exception as e:
        print(f"Stable Diffusion Error: {e}")
        return None

def get_slide_layouts(template_name):
    """
    Extracts all available slide layout names from a given PowerPoint template.

    Args:
        template_name (str): Name of the template file (without .pptx extension).

    Returns:
        list: A list of available slide layouts formatted for the LLM.
    """
    template_path = os.path.join(TEMPLATE_DIR, f"{template_name}.pptx")

    # ✅ Check if the template file exists
    if not os.path.exists(template_path):
        print(f"❌ Template '{template_name}' not found in {TEMPLATE_DIR}")
        return []

    # ✅ Open the PowerPoint template
    prs = Presentation(template_path)

    # ✅ Extract all slide layout names
    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        # Some layouts may not have names, assign generic numbering
        layout_name = layout.name if layout.name else f"Layout {i+1}"
        layouts.append({"slide_number": i, "layout_name": layout_name})

    return layouts


class SlideGenerator(BaseApp):
    def __init__(self, model_manager, app_name="Slide Generator", app_slogan="📊 Create AI-Powered Presentations!"):
        super().__init__(app_name)
        self.app_slogan = app_slogan
        self.model_manager = model_manager
        if "history" not in st.session_state:
            st.session_state.history = []

        # Initialize session state variables
        if "user_info" not in st.session_state:
            st.session_state.user_info = {
                "num_slides": 10,
                "presentation_request": "",
                "document_content": "",
            }
        if "generate_clicked" not in st.session_state:
            st.session_state.generate_clicked = False

    def welcome_screen(self):
        super().welcome_screen(self.app_slogan)

    def display_side(self):
        """Sidebar for user to specify slide preferences."""
        with st.sidebar:
            st.header("⚙️ Slide Settings")

            #  Store selected template in session state
            st.session_state.selected_template = st.selectbox(
                "🎨 Select a PowerPoint Template", ["Upload Custom"] + available_templates, index=0
            )

            if st.session_state.selected_template == "Upload Custom":
                uploaded_file = st.file_uploader("📂 Upload PowerPoint Template (.pptx)", type=["pptx"])
                if uploaded_file:
                    #  Store uploaded file in session state (not saving to disk yet)
                    st.session_state.uploaded_template = uploaded_file

            #  Ensure final_template is always a valid string
            if "uploaded_template" in st.session_state and st.session_state.uploaded_template:
                st.session_state.final_template = "Custom Upload"  #  Avoid setting it to None
            else:
                st.session_state.final_template = os.path.join(TEMPLATE_DIR, f"{st.session_state.selected_template}.pptx")

            #  Dropdown for number of slides
            st.session_state.num_slides = st.selectbox("📊 Number of Slides", [5, 10, 15, 20], index=1)

            #  Text input for user description
            st.session_state.presentation_request = st.text_area(
                "📌 Describe your presentation",
                placeholder="E.g., 'A professional presentation on AI in business with trends & ethical concerns.'"
            )

            if st.button("🎨 Generate Slides"):
                st.session_state.generate_clicked = True

    def generate_prompt(self):
        """Generate a structured prompt for AI."""

        # 🔥 Ensure final_template is a valid string before using os.path.basename
        template_name = "Custom Upload" if st.session_state.final_template == "Custom Upload" else os.path.basename(st.session_state.final_template).replace(".pptx", "")

        # 🔥 Get slide layouts based on the selected or uploaded template
        slide_layouts = get_slide_layouts(template_name)

        # 🔥 Format slide layouts for LLM understanding
        layout_text = "\n".join([f"- {layout['layout_name']} (SlideType: {layout['slide_number']})" for layout in slide_layouts])

        return ChatPromptTemplate.from_template("""
        You are an AI that generates structured PowerPoint slide content.
        When needed, suggest a detailed image description that can be used to generate an image and add it to the slide. The number of images should not exceed half the number of slides.
        If the slide does not require an image, return '#Image: None'.
        Follow this exact format but for the number of slides given

        #Title: TITLE OF THE PRESENTATION
        #Subtitle: A brief one-line description of the presentation

        #Slide: 1
        #Header: Table of Contents
        #Contents: This presentation includes:
        - Main topic 1
        - Main topic 2
        - Main topic 3
        #Image: None
        #Slidetype: CHOSEN SLIDE NUMBER

        #Slide: 2
        #Header: TITLE OF SLIDE
        #Content: CONTENT OF THE SLIDE
        #Image: SUGGESTED IMAGE description OR None
        #Slidetype: CHOSEN SLIDE NUMBER

        #Slide: 3
        #Header: TITLE OF SLIDE
        #Contents: Key points:
        - Point 1
          - Sub-point A
          - Sub-point B
        - Point 2
        #Image: SUGGESTED IMAGE description OR None
        #Slidetype: CHOSEN SLIDE NUMBER

        Ensure that all slides follow this format exactly. Do not add extra slides.

        number of slides {num_slides}
        subject {presentation_request}
        available slide layouts:
        {layout_text}
        """).format(num_slides=st.session_state.num_slides, presentation_request=st.session_state.presentation_request, layout_text=layout_text)


    def generate_ppt(self, ai_response):
        """Creates a PowerPoint presentation using AI response."""
        try:
            content = ai_response.content
            st.session_state.history.append(content)

            lines = content.split("\n")

            #  Extract all slide numbers from the LLM response
            slide_numbers = []
            slide_data = {}  # Store slide data in a structured format

            for i, line in enumerate(lines):
                line = line.strip()
                if line.startswith("#Slide:"):
                    slide_num = int(line.replace("#Slide:", "").strip())
                    slide_numbers.append(slide_num)
                    slide_data[slide_num] = {"title": "", "content": [], "image": None, "slidetype": None, "is_bullet_points": False}

            #  Ensure the uploaded template is saved before loading
            temp_template_path = os.path.join(TEMPLATE_DIR, "temp_uploaded.pptx")
            if "uploaded_template" in st.session_state and st.session_state.uploaded_template:
                with open(temp_template_path, "wb") as f:
                    f.write(st.session_state.uploaded_template.getbuffer())  # Save session file to disk

                ppt_template = temp_template_path  # Use temporary uploaded template
                theme = "Custom Upload"  #  Set theme correctly
            else:
                ppt_template = st.session_state.final_template  # Use predefined template
                theme = st.session_state.selected_template  #  Ensure theme is defined

            #  Ensure template file exists before proceeding
            if not os.path.exists(ppt_template):
                raise FileNotFoundError(f"❌ Template '{ppt_template}' not found!")

            prs = Presentation(ppt_template)

            #  If a custom template is used, remove extra slides and keep only the first one
            if theme not in available_templates and len(prs.slides) > 1:
                while len(prs.slides) > 1:
                    rId = prs.slides._sldIdLst[1].rId  # Get slide reference ID
                    prs.part.drop_rel(rId)  # Remove the reference
                    del prs.slides._sldIdLst[1]  # Delete slide from slide list

            #  Process Lines & Store Slide Data
            current_slide = None
            for line in lines:
                line = line.strip()

                if line.startswith("#Slide:"):
                    current_slide = int(line.replace("#Slide:", "").strip())
                    continue

                elif current_slide is not None:
                    if line.startswith("#Title:"):
                        slide_data[current_slide]["title"] = line.replace("#Title:", "").strip()

                    elif line.startswith("#Subtitle:"):
                        slide_data[current_slide]["subtitle"] = line.replace("#Subtitle:", "").strip()

                    elif line.startswith("#Header:"):
                        slide_data[current_slide]["title"] = line.replace("#Header:", "").strip()

                    elif line.startswith("#Contents:"):
                        slide_data[current_slide]["is_bullet_points"] = True
                        slide_data[current_slide]["content"].append(line.replace("#Contents:", "").strip())

                    elif line.startswith("#Content:"):
                        slide_data[current_slide]["is_bullet_points"] = False
                        slide_data[current_slide]["content"].append(line.replace("#Content:", "").strip())

                    elif line.startswith("- "):
                        slide_data[current_slide]["content"].append(line.replace("- ", "").strip())

                    elif line.startswith("#Image:"):
                        image_suggestion = line.replace("#Image:", "").strip()
                        slide_data[current_slide]["image"] = None if image_suggestion.lower() == "none" else image_suggestion

                    elif line.startswith("#Slidetype:"):
                        slide_data[current_slide]["slidetype"] = int(line.replace("#Slidetype:", "").strip())

            for slide_num in sorted(slide_numbers):
                slide_info = slide_data[slide_num]

                if slide_num == 1:
                    #  Title Slide
                    first_slide = prs.slides[0]
                    first_slide.shapes.title.text = slide_info["title"]

                    subtitle_placeholder = None
                    for shape in first_slide.shapes:
                        if shape.has_text_frame and shape.text_frame.text == "":
                            subtitle_placeholder = shape.text_frame
                            break

                    # If subtitle is available, add either subtitle or content
                    if subtitle_placeholder:
                        if slide_info.get("subtitle"):
                            subtitle_placeholder.text = slide_info["subtitle"]
                        elif slide_info["content"]:  #  Move content to subtitle if no subtitle is found
                            subtitle_placeholder.text = "\n".join(slide_info["content"])

                else:
                    #  Content Slides
                    slide = prs.slides.add_slide(prs.slide_layouts[int(slide_info["slidetype"])])

                    slide.shapes.title.text = slide_info["title"]

                    #  Find Correct Text Placeholder
                    text_placeholder = find_text_placeholder(slide)

                    #  Insert Content
                    if text_placeholder:
                        if slide_info["is_bullet_points"]:
                            for i, point in enumerate(slide_info["content"]):
                                p = text_placeholder.add_paragraph()
                                p.text = point.strip()
                                p.font.size = Pt(24) if i == 0 else Pt(18)
                                if i > 0:
                                    p.level = 1
                        else:
                            text_placeholder.text = "\n".join(slide_info["content"])
                            text_placeholder.paragraphs[0].font.size = Pt(24)

                    else:
                        #  If no content placeholder exists, move content into subtitle if available
                        subtitle_placeholder = None
                        for shape in slide.shapes:
                            if shape.has_text_frame and shape.text_frame.text == "":
                                subtitle_placeholder = shape.text_frame
                                break

                        if subtitle_placeholder:
                            subtitle_placeholder.text = "\n".join(slide_info["content"])


                    #  Generate & Insert Image
                    if slide_info["image"]:
                        prompt = f"You are an AI designed to make pictures for PowerPoint presentations. Make a picture about {slide_info['image']}."
                        generated_image = generate_image_stable_diffusion(prompt)

                        if generated_image:
                            img_path = os.path.join("./images", f"{slide_num}.jpg")
                            generated_image.save(img_path, format="JPEG")

                            image_placeholder = find_image_placeholder(slide)

                            if image_placeholder:
                                image_placeholder.insert_picture(img_path)
                            else:
                                slide.shapes.add_picture(img_path, Inches(1), Inches(2), Inches(6), Inches(4))

                            os.remove(img_path)

            #  Save PowerPoint
            ppt_filename = "Generated_Presentation.pptx"
            prs.save(ppt_filename)

            #  Delete the temporary uploaded template after generating the PPT
            if os.path.exists(temp_template_path):
                try:
                    os.remove(temp_template_path)
                    print(f"🗑️ Deleted temporary uploaded template: {temp_template_path}")
                except Exception as e:
                    print(f"⚠️ Could not delete temporary uploaded template: {e}")

            return ppt_filename if os.path.exists(ppt_filename) else None

        except Exception as e:
            st.error(f"Error generating PowerPoint: {e}")
            return None


    def handle_input(self):
        """Handles user input and generates a PowerPoint presentation using ModelManager."""

        if st.session_state.generate_clicked:
            if not st.session_state.presentation_request.strip():
                st.warning("Please describe the presentation you want.")
                return

            #  Prepare input for AI model
            user_input = {
                "num_slides": st.session_state.num_slides,
                "presentation_request": st.session_state.presentation_request
            }

            #  Generate AI response using ModelManager
            prompt_template = self.generate_prompt()
            response = self.model_manager.generate(prompt_template, user_input)

            #  Store response in session history
            if response:
                st.session_state.history.append(response.content)

            #  Generate PowerPoint
            ppt_filename = self.generate_ppt(response)
            if ppt_filename:
                st.session_state.ppt_filename = ppt_filename
                st.session_state.generate_clicked = False  # Prevent re-triggering
                st.session_state.show_edit_button = True  # Ensure "Make Edits" remains visible

        #  If edit mode is enabled, show the edit interface
        if st.session_state.get("edit_mode", False):
            user_edits = st.text_area("💬 Provide feedback for edits")

            if st.button("🔄 Submit Edits"):
                if st.session_state.history:
                    last_response = st.session_state.history[-1]
                    new_request = f"{last_response}\n\n# User Edits:\n{user_edits}\n\nEnsure the revised presentation follows the exact previous format."

                    #  Send updated request to AI using ModelManager
                    response = self.model_manager.generate(self.generate_prompt(), {"presentation_request": new_request})

                    #  Store new response in session history
                    if response:
                        st.session_state.history.append(response.content)

                    #  Generate new PPT with edits
                    ppt_filename = self.generate_ppt(response)

                    if ppt_filename:
                        st.session_state.ppt_filename = ppt_filename
                        st.session_state.edit_mode = False  # Exit edit mode
                        st.session_state.show_edit_button = True  # Ensure "Make Edits" remains visible

        #  Ensure "Make Edits" and download button are always visible if a presentation exists
        if "ppt_filename" in st.session_state and st.session_state.ppt_filename:

            with open(st.session_state.ppt_filename, "rb") as f:
                st.download_button("📥 Download Presentation", f, st.session_state.ppt_filename,
                                   "application/vnd.openxmlformats-officedocument.presentationml.presentation")

            #  Only show "Make Edits" if not in edit mode
            if st.session_state.get("show_edit_button", False) and not st.session_state.get("edit_mode", False):
                if st.button("📝 Make Edits"):
                    st.session_state.edit_mode = True  # Set immediately
                    st.rerun()  #  Force UI refresh to reflect changes immediately



# Run the app
if __name__ == "__main__":
    model_manager = ModelManager(PRIMARY_MODEL, SECONDARY_MODEL)
    app = SlideGenerator(model_manager)
    app.welcome_screen()
    app.display_side()
    app.handle_input()
