import streamlit as st
import requests
import os
import random
import time  
from dotenv import load_dotenv
from BaseApp import BaseApp  
from langchain_openai import ChatOpenAI
from pptx import Presentation
from pptx.util import Pt, Inches
import base64
from io import BytesIO
from PIL import Image

# ---- NEW IMPORT for token counting
try:
    import tiktoken
except ImportError:
    tiktoken = None


# Load environment variables
load_dotenv()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
API_KEY = os.getenv("SEGMIND_API_KEY")

# Ensure API key is available
if not OPENAI_API_KEY:
    st.error("🚨 OpenAI API key is missing! Please check your .env file or system environment variables.")
    st.stop()

# ✅ Initialize session variables if not present
if "history" not in st.session_state:
    st.session_state.history = []

# For cost tracking
if "img_cost" not in st.session_state:
    st.session_state.img_cost = 0.0

if "cost" not in st.session_state:
    st.session_state.cost = 0.0

if "total_cost" not in st.session_state:
    st.session_state.total_cost = 0.0

# GPT-4o-mini cost rates in USD per token
GPT4O_MINI_INPUT_COST = 0.00000015  # $0.150 / 1,000,000
GPT4O_MINI_OUTPUT_COST = 0.00000060 # $0.600 / 1,000,000


# ------------------ HELPER: Token Counting ------------------
def approximate_token_count(text_or_messages, model="gpt-3.5-turbo"):
    """
    Approximates the number of tokens for a string or list of messages using tiktoken.
    If tiktoken is unavailable, returns 0.
    
    Args:
        text_or_messages: Either a single string or a list of dict messages like:
            [{"role": "system", "content": "..."}, {"role": "user", "content": "..."}]
        model (str): Model name for the tiktoken encoder selection.

    Returns:
        int: Approximate number of tokens.
    """
    if not tiktoken:
        # If we can't import tiktoken, just return 0
        return 0

    # For Azure or unknown model naming you might need a fallback. 
    # We'll just try to get gpt-3.5 or gpt-4 default encoders
    try:
        enc = tiktoken.encoding_for_model(model)
    except:
        enc = tiktoken.get_encoding("cl100k_base")

    # If it's a list of messages, join them in the typical ChatOpenAI format
    if isinstance(text_or_messages, list):
        # We'll simulate approximate chat format
        # There's an official formula from OpenAI but let's do a simple approach:
        text_to_encode = ""
        for msg in text_or_messages:
            role = msg.get("role", "")
            content = msg.get("content", "")
            text_to_encode += f"{role}: {content}\n"
    else:
        # It's just a single string
        text_to_encode = text_or_messages

    tokens = enc.encode(text_to_encode)
    return len(tokens)
# -----------------------------------------------------------


# 🔥 Function to find the correct text placeholder
def find_text_placeholder(slide):
    """Finds the first available text placeholder in a slide."""
    for shape in slide.shapes:
        if shape.is_placeholder and ("Content Placeholder" in shape.name or "Text Placeholder" in shape.name):
            return shape.text_frame
    return None  # No text placeholder found

# 🔥 Function to find the correct image placeholder
def find_image_placeholder(slide):
    """Finds the first available image placeholder in a slide."""
    for shape in slide.shapes:
        if shape.is_placeholder and "Picture Placeholder" in shape.name:
            return shape
    return None  # No image placeholder found


def generate_image_stable_diffusion(prompt, negative_prompt="low quality, blurry", width=904, height=1260):
    """
    Generates an image using Stable Diffusion with the specified prompt and dimensions.
    Times the request duration to calculate cost at $0.001 per second.

    Args:
        prompt (str): The description of the image to generate.
        negative_prompt (str): Elements to avoid in the image.
        width (int): Width of the generated image in pixels.
        height (int): Height of the generated image in pixels.

    Returns:
        Image or None: The generated image or None if failed.
    """
    url = "https://api.segmind.com/v1/stable-diffusion-3.5-turbo-txt2img"
    headers = {"x-api-key": API_KEY}

    random_seed = random.randint(0, 99999999999)
    data = {
        "prompt": prompt,
        "negative_prompt": negative_prompt,
        "steps": 4,
        "guidance_scale": 1,
        "seed": random_seed,
        "sampler": "dpmpp_2m",
        "scheduler": "sgm_uniform",
        "width": width,
        "height": height,
        "aspect_ratio": "custom",
        "batch_size": 1,
        "image_format": "jpeg",
        "image_quality": 95,
        "base64": True
    }

    # 🔥 Begin timing for cost calculation
    start_time = time.time()
    try:
        response = requests.post(url, json=data, headers=headers)
        response.raise_for_status()
        response_data = response.json()

        image_base64 = response_data.get("image", "")
        if image_base64:
            img_data = base64.b64decode(image_base64)
            img = Image.open(BytesIO(img_data))
            return img
        else:
            return None
    except Exception as e:
        print(f"Stable Diffusion Error: {e}")
        return None
    finally:
        # 🔥 End timing and increment image cost
        end_time = time.time()
        duration = end_time - start_time  # seconds
        # cost at $0.001 per second
        cost_for_this_call = duration * 0.001
        print(f"[Image] Duration={duration:.2f}s  => cost=${cost_for_this_call}")
        st.session_state.img_cost += cost_for_this_call


# 🔥 Define Templates Path
TEMPLATE_DIR = "./Pres_templates"
available_templates = [
    "Celestial", "Circuit", "Ion", "Mesh", "Retrospect", 
    "Atlas", "Wisp", "Gallery", "Madison", "Organic"
]

class SlideGenerator(BaseApp):
    def __init__(self, app_name="Slide Generator", app_slogan="📊 Create AI-Powered Presentations!"):
        super().__init__(app_name)
        self.app_slogan = app_slogan
        if "history" not in st.session_state:
            st.session_state.history = []

    def welcome_screen(self):
        super().welcome_screen(self.app_slogan)

    def display_side(self):
        """Sidebar for user to specify slide preferences."""
        with st.sidebar:
            st.header("⚙️ Slide Settings")

            # 🔥 Store selected template in session state & theme name
            st.session_state.theme = st.selectbox(
                "🎨 Select a PowerPoint Template", available_templates, index=0
            )

            # 🔥 Set the template path based on selection
            st.session_state.final_template = os.path.join(TEMPLATE_DIR, f"{st.session_state.theme}.pptx")

            # 🔥 Dropdown for number of slides
            st.session_state.num_slides = st.selectbox("📊 Number of Slides", [5, 10, 15, 20], index=1)

            # 🔥 Text input for user description
            st.session_state.presentation_request = st.text_area(
                "📌 Describe your presentation",
                placeholder="E.g., 'A professional presentation on AI in business with trends & ethical concerns.'"
            )

            if st.button("🎨 Generate Slides"):
                st.session_state.generate_clicked = True

    def build_conversation(self):
        """Create structured input for the AI model."""
        return [
            {
                "role": "system",
                "content": (
                    "You are an AI that generates structured PowerPoint slide content. "
                    "When needed, suggest a detailed image description that can be used to generate an image and add it to the slide. "
                    "The number of images should not exceed half the number of slides. "
                    "If the slide does not require an image, return '#Image: None'."

                    "\n\nFollow this exact format but for the number of slides given:"
                    "\n\n#Slide: 1"
                    "\n#Title: TITLE OF THE PRESENTATION"
                    "\n#Subtitle: A brief one-line description of the presentation"

                    "\n\n#Slide: 2"
                    "\n#Header: Table of Contents"
                    "\n#Contents: This presentation includes:"
                    "\n- Main topic 1"
                    "\n- Main topic 2"
                    "\n- Main topic 3"
                    "\n#Image: None"

                    "\n\n#Slide: 3"
                    "\n#Header: TITLE OF SLIDE"
                    "\n#Content: CONTENT OF THE SLIDE"
                    "\n#Image: SUGGESTED IMAGE description OR None"

                    "\n\n#Slide: 4"
                    "\n#Header: TITLE OF SLIDE"
                    "\n#Contents: Key points:"
                    "\n- Point 1"
                    "\n  - Sub-point A"
                    "\n  - Sub-point B"
                    "\n- Point 2"
                    "\n#Image: SUGGESTED IMAGE description OR None"

                    "\n\nEnsure all slides follow this format exactly. Do not add extra slides."
                )
            },
            {
                "role": "user",
                "content": (
                    f"Generate {st.session_state.num_slides} slides based on this request:\n\n"
                    f"{st.session_state.presentation_request}"
                )
            }
        ]

    def generate_ppt(self, ai_response):
        """Creates a PowerPoint presentation using AI response."""
        try:
            content = ai_response.content
            st.session_state.history.append(content)

            lines = content.split("\n")

            theme = st.session_state.theme  # Directly use selected theme
            title = None
            subtitle = None

            slide_numbers = []
            slide_data = {}  # Store slide data in a structured format

            # 🔥 Extract Title and Subtitle from lines
            for i, line in enumerate(lines):
                line = line.strip()
                if line.startswith("#Slide:"):
                    slide_num = int(line.replace("#Slide:", "").strip())
                    slide_numbers.append(slide_num)
                    slide_data[slide_num] = {
                        "title": "", 
                        "content": [], 
                        "image": None, 
                        "is_bullet_points": False
                    }
                elif line.startswith("#Title:"):
                    title = line.replace("#Title:", "").strip()
                elif line.startswith("#Subtitle:"):
                    subtitle = line.replace("#Subtitle:", "").strip()

            # 🔥 Load PowerPoint Template
            ppt_template = os.path.join(TEMPLATE_DIR, f"{theme}.pptx")
            prs = Presentation(ppt_template) if os.path.exists(ppt_template) else Presentation()

            # 🔥 Modify First Slide Instead of Adding a New One
            first_slide = prs.slides[0]  # Get the existing first slide
            first_slide.shapes.title.text = title  # Set the title

            # Modify subtitle if there's a placeholder
            for shape in first_slide.shapes:
                if shape.has_text_frame and shape.text_frame.text == "":
                    shape.text_frame.text = subtitle
                    break

            # 🔥 Process Lines & Store Slide Data
            current_slide = None
            for line in lines:
                line = line.strip()

                if line.startswith("#Slide:"):
                    current_slide = int(line.replace("#Slide:", "").strip())
                    continue

                elif current_slide is not None:
                    if line.startswith("#Header:"):
                        slide_data[current_slide]["title"] = line.replace("#Header:", "").strip()

                    elif line.startswith("#Contents:"):
                        slide_data[current_slide]["is_bullet_points"] = True
                        slide_data[current_slide]["content"].append(
                            line.replace("#Contents:", "").strip()
                        )

                    elif line.startswith("#Content:"):
                        slide_data[current_slide]["is_bullet_points"] = False
                        slide_data[current_slide]["content"].append(
                            line.replace("#Content:", "").strip()
                        )

                    elif line.startswith("- "):
                        slide_data[current_slide]["content"].append(
                            line.replace("- ", "").strip()
                        )

                    elif line.startswith("#Image:"):
                        image_suggestion = line.replace("#Image:", "").strip()
                        slide_data[current_slide]["image"] = (
                            None if image_suggestion.lower() == "none" else image_suggestion
                        )

            # 🔥 Build Slides
            valid_slides = []
            for slide_num in sorted(slide_numbers):
                slide_info = slide_data[slide_num]

                # ✅ Skip Adding Slide if There's No Title or Content
                if not slide_info["title"] and not slide_info["content"]:
                    continue

                # 🔥 Choose Slide Layout
                slide_layout_index = 8 if slide_info["image"] else 1
                slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
                slide.shapes.title.text = slide_info["title"]
                valid_slides.append(slide)

                # 🔥 Find Correct Text Placeholder
                text_placeholder = find_text_placeholder(slide)

                # 🔥 Insert Content
                if text_placeholder:
                    if slide_info["is_bullet_points"]:
                        for i, point in enumerate(slide_info["content"]):
                            p = text_placeholder.add_paragraph()
                            p.text = point.strip()
                            p.font.size = Pt(24) if i == 0 else Pt(18)
                            if i > 0:
                                p.level = 1
                    else:
                        text_placeholder.text = "\n".join(slide_info["content"])
                        text_placeholder.paragraphs[0].font.size = Pt(24)
                else:
                    # 🔥 If no content placeholder, try the next possible
                    subtitle_placeholder = None
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape.text_frame.text == "":
                            subtitle_placeholder = shape.text_frame
                            break

                    if subtitle_placeholder:
                        subtitle_placeholder.text = "\n".join(slide_info["content"])

                # 🔥 Generate & Insert Image if needed
                if slide_info["image"]:
                    prompt = (
                        f"You are an AI designed to make pictures for PowerPoint presentations. "
                        f"Make a picture about {slide_info['image']}."
                    )
                    generated_image = generate_image_stable_diffusion(prompt)

                    if generated_image:
                        img_path = os.path.join("./images", f"{slide_num}.jpg")
                        generated_image.save(img_path, format="JPEG")

                        image_placeholder = find_image_placeholder(slide)
                        if image_placeholder:
                            image_placeholder.insert_picture(img_path)
                        else:
                            # fallback if placeholder doesn't exist
                            slide.shapes.add_picture(img_path, Inches(1), Inches(2), Inches(6), Inches(4))

                        os.remove(img_path)

            # 🔥 Cleanup: Remove Empty Slides
            slides_to_remove = [
                s for s in prs.slides
                if not any(shape.has_text_frame and shape.text_frame.text.strip() for shape in s.shapes)
            ]
            for s in slides_to_remove:
                prs.slides._sldIdLst.remove(s._element)

            # 🔥 Save PowerPoint
            ppt_filename = "Generated_Presentation.pptx"
            prs.save(ppt_filename)
            return ppt_filename if os.path.exists(ppt_filename) else None

        except Exception as e:
            st.error(f"Error generating PowerPoint: {e}")
            return None

    def handle_input(self):
        """Handles user input and generates a PowerPoint presentation."""
        if "generate_clicked" in st.session_state and st.session_state.generate_clicked:
            if not st.session_state.presentation_request.strip():
                st.warning("Please describe the presentation you want.")
                return

            # ---------------------
            #   1) Build the conversation (system+user) 
            #      and measure input tokens
            # ---------------------
            conversation = self.build_conversation()
            input_token_count = approximate_token_count(conversation, model="gpt-4")

            # 🔥 Create the model
            model = ChatOpenAI(temperature=0.2, model="gpt-4o-mini")

            # 2) Generate response using LangChain
            response = model.invoke(conversation)

            # 3) Count the output tokens (approx)
            output_token_count = approximate_token_count(response.content, model="gpt-4")

            # 4) Calculate cost for GPT-4o-mini
            cost_in = input_token_count * GPT4O_MINI_INPUT_COST
            cost_out = output_token_count * GPT4O_MINI_OUTPUT_COST
            total_gpt_cost = cost_in + cost_out
            # Add it to the main cost
            st.session_state.cost += total_gpt_cost

            # ---------------------
            #   Continue as before
            # ---------------------
            if response:
                st.session_state.history.append(response.content)

            ppt_filename = self.generate_ppt(response)
            if ppt_filename:
                st.session_state.ppt_filename = ppt_filename
                st.session_state.generate_clicked = False
                st.session_state.show_edit_button = True

                # 🔥 Print costs and reset
                print(f"Text cost (GPT-4o-mini) so far: {st.session_state.cost}")
                print(f"Image cost (StableDiffusion) so far: {st.session_state.img_cost}")
                total = st.session_state.cost + st.session_state.img_cost
                print(f"Total cost so far: {total}")

                # Reset costs after generation
                st.session_state.cost = 0.0
                st.session_state.img_cost = 0.0
                st.session_state.total_cost = 0.0

        # ----------------------
        #   Handle "Edit Mode"
        # ----------------------
        if "edit_mode" in st.session_state and st.session_state.edit_mode:
            user_edits = st.text_area("💬 Provide feedback for edits")

            if st.button("🔄 Submit Edits"):
                if st.session_state.history:
                    last_response = st.session_state.history[-1]
                    new_request = (
                        f"{last_response}\n\n# User Edits:\n{user_edits}\n\n"
                        "Ensure the revised presentation follows the exact previous format."
                    )

                    # We have a single user message here
                    edits_conversation = [{"role": "user", "content": new_request}]
                    input_tokens_for_edits = approximate_token_count(edits_conversation, model="gpt-4o-mini")

                    # Send updated request to LLM
                    model = ChatOpenAI(temperature=0.2, model="gpt-4o-mini")
                    response = model.invoke(edits_conversation)

                    # measure output tokens
                    out_tokens_for_edits = approximate_token_count(response.content, model="gpt-4o-mini")

                    # cost for these edits
                    cost_in = input_tokens_for_edits * GPT4O_MINI_INPUT_COST
                    cost_out = out_tokens_for_edits * GPT4O_MINI_OUTPUT_COST
                    edit_cost = cost_in + cost_out
                    st.session_state.cost += edit_cost

                    if response:
                        st.session_state.history.append(response.content)

                    ppt_filename = self.generate_ppt(response)
                    if ppt_filename:
                        st.session_state.ppt_filename = ppt_filename
                        st.session_state.edit_mode = False
                        st.session_state.show_edit_button = True

                        # 🔥 Print updated costs and reset
                        print(f"Text cost (GPT-4o-mini) so far: {st.session_state.cost:.6f}")
                        print(f"Image cost (StableDiffusion) so far: {st.session_state.img_cost:.6f}")
                        total = st.session_state.cost + st.session_state.img_cost
                        print(f"Total cost so far: {total:.6f}")

                        st.session_state.cost = 0.0
                        st.session_state.img_cost = 0.0
                        st.session_state.total_cost = 0.0

        # Ensure "Make Edits" button is always visible if a presentation exists
        if "ppt_filename" in st.session_state and st.session_state.ppt_filename:
            with open(st.session_state.ppt_filename, "rb") as f:
                st.download_button(
                    "📥 Download Presentation",
                    f,
                    st.session_state.ppt_filename,
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

            # Only show "Make Edits" if we're NOT in edit mode
            if "show_edit_button" in st.session_state and st.session_state.show_edit_button:
                if not st.session_state.get("edit_mode", False):
                    if st.button("📝 Make Edits"):
                        st.session_state.edit_mode = True
                        st.rerun()


# Run the app
if __name__ == "__main__":
    app = SlideGenerator()
    app.display_side()
    app.handle_input()
